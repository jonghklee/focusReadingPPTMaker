import sys
import os
import platform
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor  # Import RGBColor for setting font color
from PyQt5.QtWidgets import QApplication, QLabel, QVBoxLayout, QWidget, QLineEdit, QPushButton, QStackedWidget, QMessageBox, QTextEdit
from PyQt5.QtCore import Qt
from openpyxl import Workbook

class SimpleGUI(QWidget):
    def __init__(self):
        super().__init__()
        self.initUI()

    def initUI(self):
        self.setFixedSize(500, 400)  # 창 크기 고정

        self.stacked_widget = QStackedWidget(self)
        self.main_page = QWidget()
        self.excel_page = ExcelUI(self)
        self.ppt_page = PptUI(self)

        self.stacked_widget.addWidget(self.main_page)
        self.stacked_widget.addWidget(self.excel_page)
        self.stacked_widget.addWidget(self.ppt_page)

        main_layout = QVBoxLayout(self.main_page)

        self.title_label = QLabel('', self)
        main_layout.addWidget(self.title_label)

        self.paragraph_label = QLabel('''
                                      <strong>포커스리딩 훈련용 ppt 제작 소프트웨어</strong><br/>
                                      키워드 포착 훈련(p.195)을 위한 ppt 제작 소프트웨어입니다.<br/><br/><br/>
                                      <strong>1. 우선 훈련에 사용할 텍스트를 선택하세요!</strong><br/>
                                      전자책을 통해 텍스트를 복사해올 것을 추천드립니다.<br/><br/><br/>
                                      ''', self)
        main_layout.addWidget(self.paragraph_label)

        self.to_excel_button = QPushButton('2. 텍스트 추출하기', self)
        self.to_excel_button.setFixedSize(450, 100)  # 버튼 크기 설정
        self.to_excel_button.clicked.connect(self.showExcelUI)
        main_layout.addWidget(self.to_excel_button, alignment=Qt.AlignCenter)

        self.to_ppt_button = QPushButton('3. PPT 생성하기', self)
        self.to_ppt_button.setFixedSize(450, 100)  # 버튼 크기 설정
        self.to_ppt_button.clicked.connect(self.showPptUI)
        main_layout.addWidget(self.to_ppt_button, alignment=Qt.AlignCenter)

        self.setLayout(QVBoxLayout())
        self.layout().addWidget(self.stacked_widget)
        self.setWindowTitle('키워드 포착 훈련 보조 도구')
        self.show()

    def showExcelUI(self):
        self.stacked_widget.setCurrentWidget(self.excel_page)

    def showPptUI(self):
        self.stacked_widget.setCurrentWidget(self.ppt_page)

class ExcelUI(QWidget):
    def __init__(self, main_gui):
        super().__init__()
        self.main_gui = main_gui
        self.initUI()

    def initUI(self):
        layout = QVBoxLayout()

        self.back_button = QPushButton('<', self)
        self.back_button.setFixedSize(30, 30)  # 작은 버튼 크기 설정
        self.back_button.clicked.connect(self.goBack)
        layout.addWidget(self.back_button, alignment=Qt.AlignTop | Qt.AlignLeft)

        self.paragraph_label = QLabel('<b>2. 고른 텍스트를 아래에 붙여넣기 하세요.<br/>텍스트를 문장별로 나누어 엑셀 파일로 저장해드립니다.</b>', self)
        layout.addWidget(self.paragraph_label)

        self.input_field = QTextEdit(self)
        self.input_field.setPlaceholderText('''여기에 전체 텍스트를 입력하세요.\n\nex)독특한 발상과 놀라운 주장으로 40여 년간 수많은 찬사와 논쟁의 중심에 선 과학 교양서의 바이블! 1976년, 처음 출간되었을 당시 과학계와 일반 대중들에게 폭발적인 반향을 불러일으키며 세기의 문제작으로 떠오른 『이기적 유전자』는 40년이라는 세월의 검증을 거치며 그 중요성과 깊이를 더욱더 확고하게 인정받았고, 25개 이상의 언어로 번역되었으며 젊은이들이 꼭 읽어야 할 과학계의 고전으로 자리 잡았다.''')
        self.input_field.setFixedSize(450, 220)  # 폼 크기 설정
        self.input_field.setLineWrapMode(QTextEdit.WidgetWidth)  # 자동 줄바꿈 설정
        self.input_field.setVerticalScrollBarPolicy(Qt.ScrollBarAlwaysOn)  # 수직 스크롤바 항상 표시 설정
        layout.addWidget(self.input_field)

        self.submit_button = QPushButton('엑셀로 저장', self)
        self.submit_button.clicked.connect(self.save_to_excel)
        layout.addWidget(self.submit_button)

        self.setLayout(layout)
        self.setWindowTitle('엑셀로 저장')

    def goBack(self):
        self.main_gui.stacked_widget.setCurrentWidget(self.main_gui.main_page)

    def save_to_excel(self):
        text = self.input_field.toPlainText()
        sentences = [sentence.strip() for sentence in re.findall(r'.*?[.!?]+[\'\"]*', text)]
        wb = Workbook()
        ws = wb.active
        ws.append(['문장', '키워드'])
        for sentence in sentences:
            ws.append([sentence, ''])
        try:
            wb.save('focus_reading.xlsx')
            QMessageBox.information(self, '저장 완료', 'focus_reading.xlsx로 저장되었습니다!')
        except Exception as e:
            QMessageBox.critical(self, '오류', '파일 저장 중 오류가 발생했습니다: ' + str(e) + '<br/>파일을 우선 닫아주세요!')

class PptUI(QWidget):
    def __init__(self, main_gui):
        super().__init__()
        self.main_gui = main_gui
        self.initUI()

    def initUI(self):
        layout = QVBoxLayout()

        self.back_button = QPushButton('<', self)
        self.back_button.setFixedSize(30, 30)  # 작은 버튼 크기 설정
        self.back_button.clicked.connect(self.goBack)
        layout.addWidget(self.back_button, alignment=Qt.AlignTop | Qt.AlignLeft)

        self.info_label = QLabel('<b>3-1. focus_reading.xlsx로 들어가서 키워드를 입력하세요.</b><br/>(chatGPT를 사용하면 훨씬 빠르게 키워드를 입력할 수 있습니다.)', self)
        layout.addWidget(self.info_label)

        self.create_ppt_button = QPushButton('엑셀파일 열기', self)
        self.create_ppt_button.clicked.connect(self.open_excel)
        layout.addWidget(self.create_ppt_button)

        self.info_label = QLabel('<b>3-2. 키워드 선택이 끝난 뒤, 파일을 저장하고 닫아주세요.</b>', self)
        layout.addWidget(self.info_label)

        self.info_label = QLabel("<b>3-3. PPT 생성 버튼을 클릭하면 ppt가 완성됩니다.</b><br/>'focus_reading.pptx'라는 이름으로 파일이 생성됩니다.<br/>ppt파일은 실행파일이 있는 폴더에 만들어집니다.", self)
        layout.addWidget(self.info_label)

        self.create_ppt_button = QPushButton('PPT 생성', self)
        self.create_ppt_button.clicked.connect(self.create_ppt)
        layout.addWidget(self.create_ppt_button)

        self.setLayout(layout)
        self.setWindowTitle('PPT 생성기')

    def goBack(self):
        self.main_gui.stacked_widget.setCurrentWidget(self.main_gui.main_page)

    def open_excel(self):

        file_path = 'focus_reading.xlsx'
        if platform.system() == 'Windows':
            os.startfile(file_path)
        elif platform.system() == 'Darwin':  # macOS
            os.system(f'open "{file_path}"')
        else:  # Linux and others
            os.system(f'xdg-open "{file_path}"')

    def create_ppt(self):
        prs = Presentation()
        
        from openpyxl import load_workbook
        wb = load_workbook(filename='focus_reading.xlsx')
        ws = wb.active
        sentence_list = [[cell.value for cell in row] for row in ws.iter_rows(min_row=2)]
        
        for sentence, keyword in sentence_list:
            sentence = str(sentence)
            keyword = str(keyword)
            start_idx = sentence.find(keyword)
            end_idx = start_idx + len(keyword)
            # 새 슬라이드 추가
            slide_layout = prs.slide_layouts[6]  # 제목 없는 빈 슬라이드 레이아웃
            slide = prs.slides.add_slide(slide_layout)

            # 텍스트 상자 추가 및 문장 삽입
            left = Inches(1)
            top = Inches(1)
            width = Inches(8)
            height = Inches(1)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True  # 줄바꿈 활성화

            p = tf.add_paragraph()
            if start_idx == -1:
              p.text = sentence
            else:
              p.text = sentence[:start_idx]  # 시작부터 키워드 시작 전까지의 텍스트

              run = p.add_run()  # 키워드를 위한 새로운 run 추가
              run.text = sentence[start_idx:end_idx]  # 키워드 텍스트부터 문장 끝까지
              run.font.color.rgb = RGBColor(255, 0, 0)  # 텍스트 색상을 빨간색으로 설정

              run2 = p.add_run()
              run2.text = sentence[end_idx:]  # 키워드 이후 텍스트
              run2.font.color.rgb = RGBColor(0, 0, 0)  # 텍스트 색상을 검은색으로 설정

            p.font.size = Pt(24)  # 전체 텍스트 글자 크기 설정
        
        # 프레젠테이션 저장
        prs.save('focus_reading.pptx')
        QMessageBox.information(self, '저장 완료', 'focus_reading.pptx로 저장되었습니다!')

if __name__ == '__main__':
    app = QApplication(sys.argv)
    main_ui = SimpleGUI()
    sys.exit(app.exec_())